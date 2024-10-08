import random
import string
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

scientific_words = [
    "Quantum", "Molecular", "Genetic", "Neural", "Atomic", "Cellular", "Gravitational",
    "Electromagnetic", "Thermodynamic", "Biochemical", "Astronomical", "Evolutionary",
    "Nanotechnology", "Bioinformatics", "Photosynthesis", "Spectroscopy", "Entropy",
    "Genome", "Algorithm", "Hypothesis", "Synthesis", "Analysis", "Catalyst", "Polymer"
]

def generate_random_filename(course_name):
    random_numbers = ''.join(random.choices(string.digits, k=5))
    return f"{course_name}{random_numbers}.pptx"

def add_random_text_box(slide):
    left = Inches(random.uniform(0, 8))
    top = Inches(random.uniform(0, 5))
    width = Inches(random.uniform(2, 4))
    height = Inches(random.uniform(0.5, 2))

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = ' '.join(random.choices(scientific_words, k=random.randint(3, 8)))

def add_random_table(slide):
    rows = random.randint(2, 5)
    cols = random.randint(2, 5)
    left = Inches(random.uniform(0, 6))
    top = Inches(random.uniform(0, 4))
    width = Inches(random.uniform(3, 6))
    height = Inches(rows * 0.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    for row in range(rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.text = random.choice(scientific_words)

def create_random_presentation(filename):
    prs = Presentation()
    for _ in range(random.randint(7, 14)): 
        slide_layout = prs.slide_layouts[random.randint(0, len(prs.slide_layouts) - 1)]
        slide = prs.slides.add_slide(slide_layout)

        for _ in range(random.randint(1, 3)):
            add_random_text_box(slide)

        if random.choice([True, False]):
            add_random_table(slide)

    prs.save(filename)

def setup_output_folder():
    output_folder = "output"
    if os.path.exists(output_folder):
        for file_name in os.listdir(output_folder):
            file_path = os.path.join(output_folder, file_name)
            if os.path.isfile(file_path):
                os.remove(file_path)
    else:
        os.makedirs(output_folder)
    return output_folder

def main(course_name):
    output_folder = setup_output_folder()

    for _ in range(10):
        filename = generate_random_filename(course_name)
        full_path = os.path.join(output_folder, filename)
        create_random_presentation(full_path)
        print(f"Created: {full_path}")

if __name__ == "__main__":
    if len(sys.argv) > 1:
        course_name = sys.argv[1]
    else:
        course_name = input("Enter the course name: ")

    main(course_name)